"""
Lecture Notes -> Handwritten Notes
------------------------------------
Run: python app.py
Then open http://localhost:5000
Requires: pip install flask openai PyMuPDF Pillow
Put Caveat-VariableFont_wght.ttf in the same folder.
"""

import os, random, math, io, base64, traceback
from flask import Flask, request, jsonify, render_template_string
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import fitz  # PyMuPDF
from openai import OpenAI
from pptx import Presentation

# ── CONFIG ────────────────────────────────────────────────────────────────────
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "YOUR_API_KEY_HERE")
FONT_PATH      = "Biro_Script_reduced.ttf"
FONT_FALLBACK  = "Caveat-VariableFont_wght.ttf"

PAGE_W, PAGE_H      = 2550, 3300   # 300dpi letter (8.5x11in)
MARGIN_LEFT         = 380          # left red margin line x
MARGIN_RIGHT        = 2170         # faded right margin line x (~380px from right edge)
LINE_SPACING        = 60           # real lined paper spacing at 300dpi
FIRST_LINE_Y        = 320
FONT_SIZE           = 72
HEADING_SIZE        = 90
SUB_SIZE            = 76
INK                 = (30, 30, 30)
PAPER_BG            = (255, 255, 255)
LINE_COLOR          = (173, 206, 225)
MARGIN_COLOR        = (205, 80,  80)
MARGIN_RIGHT_COLOR  = (225, 170, 170)
CHAR_ROTATION_RANGE = 0.2
BASELINE_NOISE      = 0.08
SIZE_VARIATION      = 0.02
SPACING_VARIATION   = 0.04
INK_VARIATION       = 22
WORD_SPACING_JITTER = 6


# ── FONT HELPERS ──────────────────────────────────────────────────────────────

# Characters known to be missing or broken in Biro Script — always use Caveat
BIRO_MISSING = set("()-+=/\\[]{}|^~`@#$%&*<>0123456789gf")

def _pick_font(char, base_size, force_fallback=False):
    """Return the best ImageFont for this character at this size."""
    use_fallback = force_fallback or (char in BIRO_MISSING)
    path = FONT_FALLBACK if use_fallback else FONT_PATH
    try:
        return ImageFont.truetype(path, base_size), use_fallback
    except Exception:
        # Try the other font before giving up
        try:
            return ImageFont.truetype(FONT_FALLBACK if not use_fallback else FONT_PATH, base_size), True
        except Exception:
            return ImageFont.load_default(), True


# ── RENDERING PIPELINE ────────────────────────────────────────────────────────

def create_paper(draw):
    y = FIRST_LINE_Y
    while y < PAGE_H:
        draw.line([(0, y), (PAGE_W, y)], fill=LINE_COLOR, width=2)
        y += LINE_SPACING
    draw.line([(MARGIN_LEFT, 0), (MARGIN_LEFT, PAGE_H)], fill=MARGIN_COLOR, width=4)
    draw.line([(MARGIN_RIGHT, 0), (MARGIN_RIGHT, PAGE_H)], fill=MARGIN_RIGHT_COLOR, width=3)


def render_char(canvas, char, cx, baseline_y, base_size,
                rotation=0.2, noise=0.08, size_var=0.02, space_var=0.04):
    """
    Render a single character so its visual baseline sits on `baseline_y`.

    Key fix: after rotation the tile is larger; we position it so the
    *bottom of the glyph* (not the top of the tile) lands on baseline_y.
    This keeps text inside the ruled lines regardless of rotation.
    """
    size_delta = int(base_size * random.uniform(-size_var, size_var))
    char_size  = max(10, base_size + size_delta)

    font, _ = _pick_font(char, char_size)

    bbox   = font.getbbox(char)
    char_w = max(1, bbox[2] - bbox[0])
    char_h = max(4, bbox[3] - bbox[1])
    # ascent = distance from top of bbox to the drawing origin
    ascent = -bbox[1]   # bbox[1] is usually negative (above origin)

    pad  = 24
    tile = Image.new("RGBA", (char_w + pad * 2, char_h + pad * 2), (0, 0, 0, 0))
    td   = ImageDraw.Draw(tile)

    v   = random.randint(-INK_VARIATION, INK_VARIATION // 2)
    ink = (max(0, min(255, INK[0] + v)),
           max(0, min(255, INK[1] + v)),
           max(0, min(255, INK[2] + v)),
           random.randint(210, 255))

    # Draw glyph inside the padded tile
    td.text((pad - bbox[0], pad - bbox[1]), char, font=font, fill=ink)

    # Rotate the tile (expand=True keeps the full rotated image)
    angle      = random.uniform(-rotation, rotation)
    rotated    = tile.rotate(math.degrees(angle) if abs(angle) < 3 else angle,
                              expand=True, resample=Image.BICUBIC)

    rw, rh = rotated.size

    # Baseline sits at (pad + ascent + char_h) from the tile top before rotation.
    # After rotation with expand=True the centre of the tile is preserved.
    # We approximate: place the tile so the bottom of the original glyph area
    # aligns to baseline_y, then add a tiny y-noise that is MUCH smaller than
    # half the line spacing so the text stays inside the lines.
    max_y_jitter = LINE_SPACING * noise          # e.g. 0.08 * 60 = ~5px
    y_noise      = random.uniform(-max_y_jitter, max_y_jitter)

    # The glyph bottom (before rotation) is at: pad + char_h inside the tile.
    # After expand-rotate the tile grew; the original centre is at (rw/2, rh/2).
    # We want the glyph bottom to sit at baseline_y:
    #   tile_centre_y_before = pad + (-bbox[1]) + char_h - char_h/2  (approx)
    # Simpler and robust: anchor to cap-height, not full descent.
    cap_y_in_tile = pad        # top of the padded glyph ~ cap top
    glyph_bottom_in_tile = pad + char_h

    # y position of tile so that glyph_bottom_in_tile aligns to baseline_y
    paste_y = int(baseline_y - glyph_bottom_in_tile + y_noise)

    # x: start at cx, shift back by pad so glyph starts at cx
    paste_x = int(cx) - pad

    canvas.paste(rotated, (paste_x, paste_y), rotated)

    advance = char_w + int(char_w * random.uniform(-space_var, space_var))
    return max(4, advance)


def render_text_line(canvas, text, x_start, y_baseline, base_size,
                     rotation=0.2, noise=0.08, size_var=0.02, space_var=0.04):
    x = x_start
    for char in text:
        if char == " ":
            x += int(base_size * 0.28) + random.randint(-WORD_SPACING_JITTER, WORD_SPACING_JITTER)
            continue
        x += render_char(canvas, char, x, y_baseline, base_size,
                         rotation=rotation, noise=noise,
                         size_var=size_var, space_var=space_var)
        if x > MARGIN_RIGHT - 30:
            break


def draw_underline(draw_canvas, x_start, y, text, size):
    try:
        font = ImageFont.truetype(FONT_PATH, size)
    except Exception:
        return
    bbox   = font.getbbox(text)
    text_w = bbox[2] - bbox[0]
    uy     = y + (bbox[3] - bbox[1]) + 2
    px     = x_start
    while px < x_start + text_w:
        seg_len = random.randint(6, 14)
        end_x   = min(px + seg_len, x_start + text_w)
        y1 = uy + random.uniform(-1.5, 1.5)
        y2 = uy + random.uniform(-1.5, 1.5)
        v  = random.randint(-15, 5)
        draw_canvas.line([(px, y1), (end_x, y2)],
                         fill=(max(0, INK[0] + v), max(0, INK[1] + v), max(0, INK[2] + v)),
                         width=random.randint(1, 2))
        px = end_x + random.randint(0, 2)


def measure_text_width(text, font_size):
    """Measure how wide a string will render at a given font size."""
    total = 0
    for char in text:
        if char == " ":
            total += int(font_size * 0.28)
        else:
            font, _ = _pick_font(char, font_size)
            bbox     = font.getbbox(char)
            total   += max(4, bbox[2] - bbox[0])
    return total


def wrap_text(text, x_start, font_size, max_x):
    """Break text into lines that fit within max_x, returning list of strings."""
    max_width = max_x - x_start
    words     = text.split(" ")
    lines     = []
    current   = ""
    for word in words:
        test = (current + " " + word).strip()
        if measure_text_width(test, font_size) <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines if lines else [text]


def render_page(lines, rotation=0.2, noise=0.08, size_var=0.02, space_var=0.04):
    img    = Image.new("RGB",  (PAGE_W, PAGE_H), PAPER_BG)
    canvas = Image.new("RGBA", (PAGE_W, PAGE_H), (0, 0, 0, 0))
    draw_bg     = ImageDraw.Draw(img)
    draw_canvas = ImageDraw.Draw(canvas)
    create_paper(draw_bg)

    margin = MARGIN_LEFT + 30

    # Pull the first # heading out and render it in the title area above the rules
    title_line = None
    start_i    = 0
    for idx, line in enumerate(lines):
        if line.strip():
            if line.startswith("# "):
                title_line = line
                start_i    = idx + 1
            break

    if title_line:
        title_text = title_line[2:].rstrip()
        title_y    = FIRST_LINE_Y - 20
        tx         = margin + random.randint(-4, 8)
        render_text_line(canvas, title_text, tx, title_y, HEADING_SIZE,
                         rotation=rotation, noise=noise, size_var=size_var, space_var=space_var)
        # Body starts on the second ruled line (title lives above the rules)
        y = FIRST_LINE_Y + LINE_SPACING * 2 + 4
    else:
        # No title — skip the title gap, start on the second ruled line
        y = FIRST_LINE_Y + LINE_SPACING * 2 + 4

    i = start_i
    while i < len(lines):
        raw = lines[i].rstrip()

        if raw.startswith("# "):
            text = raw[2:]
            x    = margin + random.randint(-4, 8)
            wrapped_lines = wrap_text(text, x, HEADING_SIZE, MARGIN_RIGHT - 20)
            if y + int(LINE_SPACING * 2.2) * len(wrapped_lines) > PAGE_H - 150:
                break
            for j, wrapped in enumerate(wrapped_lines):
                render_text_line(canvas, wrapped, x, y, HEADING_SIZE,
                                 rotation=rotation, noise=noise, size_var=size_var, space_var=space_var)
                y += int(LINE_SPACING * 2.2) if j == len(wrapped_lines) - 1 else int(LINE_SPACING * 1.5)

        elif raw.startswith("## "):
            text = raw[3:]
            x    = margin + random.randint(-2, 6)
            wrapped_lines = wrap_text(text, x, SUB_SIZE, MARGIN_RIGHT - 20)
            if y + int(LINE_SPACING * 1.6) * len(wrapped_lines) > PAGE_H - 150:
                break
            for j, wrapped in enumerate(wrapped_lines):
                render_text_line(canvas, wrapped, x, y, SUB_SIZE,
                                 rotation=rotation, noise=noise, size_var=size_var, space_var=space_var)
                y += int(LINE_SPACING * 1.6) if j == len(wrapped_lines) - 1 else int(LINE_SPACING * 1.1)

        elif raw.startswith("  - "):
            text = "- " + raw[4:]
            x    = margin + 120 + random.randint(-4, 6)
            cont = margin + 160
            wrapped_lines = wrap_text(text, x, FONT_SIZE - 2, MARGIN_RIGHT - 20)
            if y + LINE_SPACING * len(wrapped_lines) > PAGE_H - 150:
                break
            for j, wrapped in enumerate(wrapped_lines):
                render_text_line(canvas, wrapped, x if j == 0 else cont, y, FONT_SIZE - 2,
                                 rotation=rotation, noise=noise, size_var=size_var, space_var=space_var)
                y += LINE_SPACING

        elif raw.startswith("- "):
            text = "- " + raw[2:]
            x    = margin + 30 + random.randint(-4, 6)
            cont = margin + 60
            wrapped_lines = wrap_text(text, x, FONT_SIZE, MARGIN_RIGHT - 20)
            if y + LINE_SPACING * len(wrapped_lines) > PAGE_H - 150:
                break
            for j, wrapped in enumerate(wrapped_lines):
                render_text_line(canvas, wrapped, x if j == 0 else cont, y, FONT_SIZE,
                                 rotation=rotation, noise=noise, size_var=size_var, space_var=space_var)
                y += LINE_SPACING

        elif raw == "":
            y += int(LINE_SPACING * 0.55)

        else:
            x = margin + random.randint(-4, 10)
            wrapped_lines = wrap_text(raw, x, FONT_SIZE, MARGIN_RIGHT - 20)
            if y + LINE_SPACING * len(wrapped_lines) > PAGE_H - 150:
                break
            for wrapped in wrapped_lines:
                render_text_line(canvas, wrapped, x, y, FONT_SIZE,
                                 rotation=rotation, noise=noise, size_var=size_var, space_var=space_var)
                y += LINE_SPACING

        i += 1

    img.paste(
        Image.alpha_composite(Image.new("RGBA", (PAGE_W, PAGE_H), (0, 0, 0, 0)), canvas).convert("RGB"),
        mask=canvas.split()[3]
    )
    img = img.filter(ImageFilter.GaussianBlur(radius=1.2))
    return img, lines[i:]


def render_notes_to_b64(notes_text, messiness=0.5):
    # Keep rotation gentle so text stays inside ruled lines.
    # Even at messiness=1.0 the per-char tilt is max ~4°, which looks messy
    # but still readable and on the line.
    rotation  = messiness * 0.35    # 0 → 0°,   1 → ~0.35 rad ≈ 20° max per char — still readable
    noise     = messiness * 0.06    # baseline jitter as fraction of line spacing
    size_var  = messiness * 0.04
    space_var = messiness * 0.06

    remaining = notes_text.strip().split("\n")
    images    = []
    while remaining:
        img, remaining = render_page(remaining, rotation, noise, size_var, space_var)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        images.append(base64.b64encode(buf.getvalue()).decode())
        if not remaining:
            break
    return images


# ── GPT CALL ─────────────────────────────────────────────────────────────────

def extract_from_pdf(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = "\n\n".join(page.get_text() for page in doc)
    if not text.strip():
        raise ValueError("Could not extract text from Lecture (might be scanned/image-based)")
    return text

def extract_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)
        if parts:
            slides.append(f"[Slide {i+1}]\n" + "\n".join(parts))
    text = "\n\n".join(slides)
    if not text.strip():
        raise ValueError("Could not extract text from PPTX")
    return text

def extract_from_image(file_bytes, filename):
    client   = OpenAI(api_key=OPENAI_API_KEY)
    img_b64  = base64.b64encode(file_bytes).decode()
    ext      = filename.rsplit(".", 1)[-1].lower()
    mime     = "image/png" if ext == "png" else "image/jpeg"
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{img_b64}"}},
                {"type": "text",
                 "text": "Extract all the text content from this image exactly as it appears. "
                         "Include all text from slides, diagrams, bullet points, headings, and labels. "
                         "Output only the extracted text, no commentary."}
            ]
        }],
        max_tokens=4000
    )
    text = response.choices[0].message.content.strip()
    if not text:
        raise ValueError("Could not extract text from image")
    return text

def extract_from_upload(file_bytes, filename):
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext == "pdf":
        return extract_from_pdf(file_bytes)
    elif ext == "pptx":
        return extract_from_pptx(file_bytes)
    elif ext in ("png", "jpg", "jpeg", "webp"):
        return extract_from_image(file_bytes, filename)
    else:
        raise ValueError(f"Unsupported file type: .{ext}")

def generate_notes(raw_text, detail=0.5, custom_instructions=""):
    words = raw_text.split()
    if len(words) > 6000:
        raw_text = " ".join(words[:6000]) + "\n[truncated]"

    if detail < 0.33:
        detail_instruction = (
            "- Write HIGH-LEVEL summary notes only — key topics and main ideas, no deep explanations\n"
            "- Keep bullets short, one line max\n"
            "- Aim for brevity: a student reviewing quickly before an exam"
        )
    elif detail < 0.67:
        detail_instruction = (
            "- Include specific terms, definitions, and mechanisms with brief explanations\n"
            "- Aim for the level of detail a diligent student would write\n"
            "- Balance coverage and conciseness"
        )
    else:
        detail_instruction = (
            "- Write COMPREHENSIVE notes — include all key terms, definitions, formulas, mechanisms, examples, and exceptions\n"
            "- Explain each concept clearly, as if teaching it to someone\n"
            "- Include sub-bullets for nuance and detail wherever useful\n"
            "- Prioritize completeness over brevity"
        )

    custom_block = ""
    if custom_instructions.strip():
        custom_block = f"\nAdditional instructions from the student:\n{custom_instructions.strip()}\n"

    client   = OpenAI(api_key=OPENAI_API_KEY)
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{
            "role": "user",
            "content": f"""Rewrite these lecture slides as handwritten student notes.

Rules:
{detail_instruction}
- SKIP overview slides, learning objectives, and generic filler sentences
- Use only ASCII characters - no unicode, use -> instead of arrows
- Do NOT use ** for bold, * for italic, or any markdown formatting whatsoever
- Do NOT wrap output in code blocks or markdown fences
- Structure with these exact formats:
    # Main topic heading
    ## Sub heading
    - bullet point
      - sub bullet (2 spaces then dash)
    plain paragraph text for explanations
- Abbreviations are fine (w/, b/c, approx, etc.)
- Do NOT write a summary, conclusion, or closing paragraph at the end — stop after the last topic
{custom_block}
Lecture content:
{raw_text}"""
        }]
    )
    text = response.choices[0].message.content
    import re
    text = re.sub(r'^```[a-zA-Z]*\n?', '', text.strip())
    text = re.sub(r'\n?```$', '', text.strip())
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'^#{3,}\s*', '## ', text, flags=re.MULTILINE)

    # Strip any trailing summary/conclusion paragraph GPT sneaks in.
    # These tend to be plain paragraphs starting with summary-ish phrases
    # at the very end of the output (after the last heading/bullet block).
    summary_triggers = (
        'in summary', 'in conclusion', 'overall,', 'to summarize',
        'this outline', 'this overview', 'these notes', 'understanding',
        'in short,', 'taken together', 'together, these',
    )
    lines = text.split('\n')
    # Walk backwards, drop trailing plain-text paragraphs that look like summaries
    while lines:
        last = lines[-1].strip().lower()
        if not last:
            lines.pop()
            continue
        if (not last.startswith('#') and not last.startswith('-')
                and any(last.startswith(t) for t in summary_triggers)):
            lines.pop()
        else:
            break

    return '\n'.join(lines).strip()


# ── FLASK APP ─────────────────────────────────────────────────────────────────

app = Flask(__name__)

HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>HandyWrite — Turn Lectures into Handwritten Notes</title>
<link rel="icon" type="image/png" href="/static/favicon.png">
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Playfair+Display:ital,wght@0,400;0,700;1,400&family=DM+Mono:wght@300;400&family=DM+Sans:wght@300;400;500&display=swap" rel="stylesheet">
<style>
  *, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }

  :root {
    --navy:       #0d1b2a;
    --navy-mid:   #152336;
    --navy-light: #1e3450;
    --ink-blue:   #1a3a6b;
    --royal:      #1c3f8f;
    --accent:     #4a7fd4;
    --accent-bright: #6b9fe8;
    --cream:      #f5f0e8;
    --paper:      #faf7f0;
    --paper-line: #ddd5c2;
    --gold:       #c9a84c;
    --gold-light: #e8c97a;
    --red-margin: #c94040;
    --text-main:  #e8e0d0;
    --text-muted: #8a9ab5;
    --text-dim:   #4a5a72;
    --success:    #4a9e72;
    --error:      #c94040;
    --glow:       rgba(74, 127, 212, 0.15);
    --shadow-deep: 0 20px 60px rgba(0,0,0,0.5);
  }

  html { scroll-behavior: smooth; }

  body {
    background: var(--navy);
    font-family: 'DM Sans', sans-serif;
    color: var(--text-main);
    min-height: 100vh;
    display: flex;
    flex-direction: column;
    align-items: center;
    padding: 0 0 80px;
    position: relative;
    overflow-x: hidden;
  }

  /* ── Ruled paper background lines ── */
  body::before {
    content: '';
    position: fixed;
    inset: 0;
    background-image:
      repeating-linear-gradient(
        to bottom,
        transparent 0px,
        transparent 31px,
        rgba(74, 127, 212, 0.04) 31px,
        rgba(74, 127, 212, 0.04) 32px
      );
    pointer-events: none;
    z-index: 0;
  }

  /* ── Subtle radial glow behind hero ── */
  body::after {
    content: '';
    position: fixed;
    top: -20%;
    left: 50%;
    transform: translateX(-50%);
    width: 800px;
    height: 500px;
    background: radial-gradient(ellipse, rgba(28,63,143,0.25) 0%, transparent 70%);
    pointer-events: none;
    z-index: 0;
  }

  /* ── HERO ── */
  .hero {
    position: relative;
    z-index: 1;
    width: 100%;
    display: flex;
    flex-direction: column;
    align-items: center;
    padding: 56px 24px 40px;
    text-align: center;
  }

  .logo-mark {
    width: 64px;
    height: 64px;
    margin-bottom: 20px;
    border-radius: 50%;
    border: 2px solid rgba(74,127,212,0.4);
    padding: 4px;
    box-shadow: 0 0 24px rgba(74,127,212,0.2), 0 0 48px rgba(74,127,212,0.08);
    animation: float 4s ease-in-out infinite;
  }

  @keyframes float {
    0%, 100% { transform: translateY(0px); }
    50%       { transform: translateY(-5px); }
  }

  .hero h1 {
    font-family: 'Playfair Display', Georgia, serif;
    font-size: clamp(2rem, 5vw, 3rem);
    font-weight: 700;
    letter-spacing: -0.02em;
    color: var(--cream);
    line-height: 1.1;
  }

  .hero h1 em {
    font-style: italic;
    color: var(--accent-bright);
  }

  .hero-sub {
    margin-top: 10px;
    font-size: 0.95rem;
    font-weight: 300;
    color: var(--text-muted);
    letter-spacing: 0.02em;
    max-width: 380px;
  }

  /* ── PILL BADGES ── */
  .badges {
    display: flex;
    gap: 8px;
    margin-top: 18px;
    flex-wrap: wrap;
    justify-content: center;
  }

  .badge {
    padding: 4px 12px;
    border-radius: 999px;
    border: 1px solid rgba(74,127,212,0.3);
    background: rgba(74,127,212,0.08);
    font-size: 0.72rem;
    font-family: 'DM Mono', monospace;
    color: var(--accent);
    letter-spacing: 0.05em;
    text-transform: uppercase;
  }

  /* ── MAIN CARD ── */
  .card {
    position: relative;
    z-index: 1;
    background: var(--navy-mid);
    border: 1px solid rgba(74,127,212,0.2);
    border-radius: 16px;
    padding: 32px;
    width: calc(100% - 32px);
    max-width: 600px;
    box-shadow: var(--shadow-deep), inset 0 1px 0 rgba(255,255,255,0.05);
    animation: slideUp 0.5s ease both;
  }

  @keyframes slideUp {
    from { opacity: 0; transform: translateY(24px); }
    to   { opacity: 1; transform: translateY(0); }
  }

  /* ── DROP ZONE ── */
  .drop-zone {
    border: 2px dashed rgba(74,127,212,0.3);
    border-radius: 12px;
    padding: 36px 24px;
    text-align: center;
    cursor: pointer;
    transition: all 0.25s cubic-bezier(0.4,0,0.2,1);
    background: rgba(13,27,42,0.6);
    position: relative;
    overflow: hidden;
  }

  .drop-zone::before {
    content: '';
    position: absolute;
    inset: 0;
    background: radial-gradient(ellipse at 50% 0%, rgba(74,127,212,0.1) 0%, transparent 60%);
    opacity: 0;
    transition: opacity 0.3s;
  }

  .drop-zone:hover::before,
  .drop-zone.drag-over::before { opacity: 1; }

  .drop-zone:hover,
  .drop-zone.drag-over {
    border-color: var(--accent);
    background: rgba(26, 58, 107, 0.3);
    transform: translateY(-1px);
    box-shadow: 0 8px 32px rgba(74,127,212,0.15);
  }

  .drop-zone input[type="file"] {
    position: absolute; inset: 0;
    opacity: 0; cursor: pointer; width: 100%; height: 100%;
  }

  .drop-icon-wrap {
    width: 52px;
    height: 52px;
    margin: 0 auto 14px;
    border-radius: 50%;
    background: linear-gradient(135deg, rgba(74,127,212,0.2), rgba(28,63,143,0.3));
    border: 1px solid rgba(74,127,212,0.3);
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 1.5rem;
    transition: transform 0.3s;
  }

  .drop-zone:hover .drop-icon-wrap { transform: scale(1.1) rotate(-5deg); }

  .drop-zone p {
    color: var(--text-muted);
    font-size: 0.88rem;
    line-height: 1.5;
  }

  .drop-zone p strong {
    color: var(--accent-bright);
    font-weight: 500;
  }

  .file-types {
    margin-top: 8px;
    font-family: 'DM Mono', monospace;
    font-size: 0.7rem;
    color: var(--text-dim);
    letter-spacing: 0.05em;
  }

  .file-name-display {
    margin-top: 12px;
    display: none;
    align-items: center;
    justify-content: center;
    gap: 6px;
    padding: 6px 14px;
    background: rgba(74,127,212,0.12);
    border: 1px solid rgba(74,127,212,0.25);
    border-radius: 6px;
    font-family: 'DM Mono', monospace;
    font-size: 0.75rem;
    color: var(--accent-bright);
  }

  .file-name-display.visible { display: flex; }
  .file-name-display .check { color: var(--success); }

  /* ── DIVIDER ── */
  .section-divider {
    display: flex;
    align-items: center;
    gap: 12px;
    margin: 22px 0 0;
  }

  .section-divider::before,
  .section-divider::after {
    content: '';
    flex: 1;
    height: 1px;
    background: rgba(74,127,212,0.15);
  }

  .section-divider span {
    font-size: 0.7rem;
    font-family: 'DM Mono', monospace;
    color: var(--text-dim);
    letter-spacing: 0.1em;
    text-transform: uppercase;
    white-space: nowrap;
  }

  /* ── INSTRUCTIONS ── */
  .instructions-wrap { margin-top: 14px; }

  .instructions-wrap label {
    display: flex;
    align-items: center;
    gap: 6px;
    font-size: 0.78rem;
    color: var(--text-muted);
    margin-bottom: 8px;
    font-family: 'DM Mono', monospace;
    letter-spacing: 0.04em;
  }

  .instructions-wrap label .tag {
    padding: 1px 7px;
    border-radius: 999px;
    background: rgba(201,168,76,0.12);
    border: 1px solid rgba(201,168,76,0.25);
    color: var(--gold);
    font-size: 0.64rem;
    letter-spacing: 0.06em;
  }

  .instructions-wrap textarea {
    width: 100%;
    padding: 12px 14px;
    border: 1px solid rgba(74,127,212,0.2);
    border-radius: 10px;
    background: rgba(13,27,42,0.6);
    font-family: 'DM Sans', sans-serif;
    font-size: 0.84rem;
    font-weight: 300;
    color: var(--text-main);
    resize: none;
    min-height: 72px;
    max-height: 140px;
    line-height: 1.6;
    transition: border-color 0.2s, box-shadow 0.2s;
    outline: none;
  }

  .instructions-wrap textarea:focus {
    border-color: rgba(74,127,212,0.5);
    box-shadow: 0 0 0 3px rgba(74,127,212,0.08);
  }

  .instructions-wrap textarea::placeholder {
    color: var(--text-dim);
    font-style: italic;
  }

  /* ── GENERATE BUTTON ── */
  .btn-wrap { margin-top: 20px; }

  .btn {
    display: flex;
    align-items: center;
    justify-content: center;
    gap: 10px;
    width: 100%;
    padding: 15px;
    background: linear-gradient(135deg, var(--royal) 0%, var(--ink-blue) 100%);
    color: var(--cream);
    border: 1px solid rgba(74,127,212,0.4);
    border-radius: 10px;
    font-family: 'Playfair Display', serif;
    font-size: 1.05rem;
    font-weight: 400;
    letter-spacing: 0.02em;
    cursor: pointer;
    transition: all 0.25s cubic-bezier(0.4,0,0.2,1);
    position: relative;
    overflow: hidden;
  }

  .btn::after {
    content: '';
    position: absolute;
    inset: 0;
    background: linear-gradient(135deg, rgba(255,255,255,0.08) 0%, transparent 50%);
    pointer-events: none;
  }

  .btn:hover:not(:disabled) {
    transform: translateY(-2px);
    box-shadow: 0 8px 28px rgba(28,63,143,0.5), 0 0 0 1px rgba(74,127,212,0.5);
  }

  .btn:active:not(:disabled) { transform: translateY(0); }

  .btn:disabled {
    background: rgba(30,52,80,0.5);
    border-color: rgba(74,127,212,0.1);
    color: var(--text-dim);
    cursor: not-allowed;
  }

  .btn-icon { font-size: 1.1rem; transition: transform 0.3s; }
  .btn:hover:not(:disabled) .btn-icon { transform: rotate(-10deg) scale(1.1); }

  /* ── STATUS ── */
  .status {
    margin-top: 12px;
    font-size: 0.82rem;
    color: var(--text-muted);
    font-style: italic;
    text-align: center;
    min-height: 18px;
    transition: color 0.2s;
  }

  .status.error { color: var(--error); font-style: normal; }

  /* ── LOADING BAR ── */
  .progress-wrap {
    margin-top: 12px;
    height: 3px;
    border-radius: 999px;
    background: rgba(74,127,212,0.12);
    overflow: hidden;
    opacity: 0;
    transition: opacity 0.3s;
  }

  .progress-wrap.visible { opacity: 1; }

  .progress-bar {
    height: 100%;
    width: 0%;
    border-radius: 999px;
    background: linear-gradient(90deg, var(--royal), var(--accent-bright));
    transition: width 0.6s cubic-bezier(0.4, 0, 0.2, 1);
    position: relative;
  }

  /* Shimmer sweep */
  .progress-bar::after {
    content: '';
    position: absolute;
    top: 0; right: 0; bottom: 0;
    width: 60px;
    background: linear-gradient(90deg, transparent, rgba(255,255,255,0.4), transparent);
    animation: shimmer 1.2s ease-in-out infinite;
  }

  @keyframes shimmer {
    0%   { transform: translateX(-60px); opacity: 0; }
    50%  { opacity: 1; }
    100% { transform: translateX(60px); opacity: 0; }
  }

  /* ── SPINNER ── */
  .spinner {
    display: inline-block;
    width: 14px; height: 14px;
    border: 2px solid rgba(255,255,255,0.15);
    border-top-color: var(--accent-bright);
    border-radius: 50%;
    animation: spin 0.75s linear infinite;
    flex-shrink: 0;
  }
  @keyframes spin { to { transform: rotate(360deg); } }

  /* ── RESULTS SECTION ── */
  .results-header {
    position: relative;
    z-index: 1;
    width: calc(100% - 32px);
    max-width: 860px;
    margin-top: 48px;
    display: flex;
    align-items: center;
    justify-content: space-between;
    gap: 16px;
    flex-wrap: wrap;
  }

  .results-title {
    font-family: 'Playfair Display', serif;
    font-size: 1.1rem;
    color: var(--cream);
    display: flex;
    align-items: center;
    gap: 8px;
  }

  .results-title .page-count {
    font-family: 'DM Mono', monospace;
    font-size: 0.72rem;
    color: var(--accent);
    background: rgba(74,127,212,0.12);
    border: 1px solid rgba(74,127,212,0.2);
    padding: 2px 8px;
    border-radius: 999px;
  }

  .dl-btn {
    padding: 9px 20px;
    background: transparent;
    border: 1px solid rgba(74,127,212,0.35);
    border-radius: 8px;
    font-family: 'DM Mono', monospace;
    font-size: 0.75rem;
    letter-spacing: 0.04em;
    cursor: pointer;
    color: var(--accent);
    text-decoration: none;
    transition: all 0.2s;
    display: flex;
    align-items: center;
    gap: 6px;
  }

  .dl-btn:hover {
    background: rgba(74,127,212,0.12);
    border-color: var(--accent);
    color: var(--accent-bright);
    transform: translateY(-1px);
  }

  /* ── PAGES OUTPUT ── */
  .pages {
    position: relative;
    z-index: 1;
    margin-top: 16px;
    width: calc(100% - 32px);
    max-width: 860px;
    display: flex;
    flex-direction: column;
    gap: 20px;
  }

  .page-wrap {
    position: relative;
    border-radius: 8px;
    overflow: hidden;
    box-shadow:
      0 2px 4px rgba(0,0,0,0.3),
      0 8px 32px rgba(0,0,0,0.4),
      0 0 0 1px rgba(255,255,255,0.04);
    animation: pageIn 0.4s ease both;
  }

  .page-wrap:nth-child(1) { animation-delay: 0.05s; }
  .page-wrap:nth-child(2) { animation-delay: 0.12s; }
  .page-wrap:nth-child(3) { animation-delay: 0.19s; }
  .page-wrap:nth-child(n+4) { animation-delay: 0.25s; }

  @keyframes pageIn {
    from { opacity: 0; transform: translateY(16px) scale(0.99); }
    to   { opacity: 1; transform: translateY(0) scale(1); }
  }

  .page-wrap img { display: block; width: 100%; height: auto; }

  .page-label {
    position: absolute;
    bottom: 12px;
    right: 14px;
    font-family: 'DM Mono', monospace;
    font-size: 0.65rem;
    color: rgba(100,120,150,0.7);
    background: rgba(13,27,42,0.7);
    backdrop-filter: blur(4px);
    padding: 3px 8px;
    border-radius: 4px;
    border: 1px solid rgba(74,127,212,0.15);
  }

  /* ── LANDING PAGE ── */
  .landing {
    position: relative;
    z-index: 1;
    min-height: 100vh;
    width: 100%;
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    padding: 60px 24px 80px;
    text-align: center;
  }

  .landing-logo {
    width: 80px;
    height: 80px;
    border-radius: 50%;
    border: 2px solid rgba(74,127,212,0.35);
    padding: 4px;
    box-shadow: 0 0 32px rgba(74,127,212,0.25), 0 0 80px rgba(74,127,212,0.08);
    animation: float 4s ease-in-out infinite;
    margin-bottom: 28px;
  }

  .landing h1 {
    font-family: 'Playfair Display', Georgia, serif;
    font-size: clamp(2.8rem, 7vw, 4.5rem);
    font-weight: 700;
    letter-spacing: -0.03em;
    color: var(--cream);
    line-height: 1.08;
  }

  .landing h1 em {
    font-style: italic;
    color: var(--accent-bright);
  }

  .landing-tagline {
    margin-top: 18px;
    font-size: clamp(1rem, 2.5vw, 1.2rem);
    font-weight: 300;
    color: var(--text-muted);
    max-width: 480px;
    line-height: 1.6;
  }

  .landing-tagline strong {
    color: var(--text-main);
    font-weight: 400;
  }

  /* Feature pills row */
  .features {
    display: flex;
    gap: 12px;
    margin-top: 40px;
    flex-wrap: wrap;
    justify-content: center;
    max-width: 560px;
  }

  .feature-card {
    background: rgba(30, 52, 80, 0.5);
    border: 1px solid rgba(74,127,212,0.18);
    border-radius: 12px;
    padding: 16px 20px;
    text-align: left;
    width: 160px;
    transition: transform 0.2s, border-color 0.2s;
  }

  .feature-card:hover {
    transform: translateY(-3px);
    border-color: rgba(74,127,212,0.4);
  }

  .feature-icon { font-size: 1.4rem; margin-bottom: 8px; display: block; }

  .feature-card h3 {
    font-family: 'Playfair Display', serif;
    font-size: 0.88rem;
    font-weight: 400;
    color: var(--cream);
    margin-bottom: 4px;
  }

  .feature-card p {
    font-size: 0.72rem;
    color: var(--text-dim);
    line-height: 1.4;
  }

  /* CTA button */
  .cta-btn {
    margin-top: 48px;
    display: inline-flex;
    align-items: center;
    gap: 10px;
    padding: 16px 36px;
    background: linear-gradient(135deg, var(--royal) 0%, var(--ink-blue) 100%);
    color: var(--cream);
    border: 1px solid rgba(74,127,212,0.5);
    border-radius: 999px;
    font-family: 'Playfair Display', serif;
    font-size: 1.1rem;
    cursor: pointer;
    text-decoration: none;
    transition: all 0.25s cubic-bezier(0.4,0,0.2,1);
    box-shadow: 0 4px 24px rgba(28,63,143,0.35);
    position: relative;
    overflow: hidden;
  }

  .cta-btn::before {
    content: '';
    position: absolute;
    inset: 0;
    background: linear-gradient(135deg, rgba(255,255,255,0.08) 0%, transparent 60%);
  }

  .cta-btn:hover {
    transform: translateY(-3px);
    box-shadow: 0 12px 40px rgba(28,63,143,0.55), 0 0 0 1px rgba(74,127,212,0.6);
  }

  .cta-arrow {
    font-size: 1rem;
    transition: transform 0.3s;
  }

  .cta-btn:hover .cta-arrow { transform: translateX(4px); }

  .scroll-hint {
    position: absolute;
    bottom: 28px;
    left: 50%;
    transform: translateX(-50%);
    display: flex;
    flex-direction: column;
    align-items: center;
    gap: 6px;
    color: var(--text-dim);
    font-family: 'DM Mono', monospace;
    font-size: 0.65rem;
    letter-spacing: 0.1em;
    text-transform: uppercase;
    animation: fadeInUp 1s ease 1s both;
  }

  .scroll-hint .chevron {
    width: 16px;
    height: 16px;
    border-right: 1.5px solid var(--text-dim);
    border-bottom: 1.5px solid var(--text-dim);
    transform: rotate(45deg);
    animation: bounce 1.5s ease-in-out infinite;
  }

  @keyframes bounce {
    0%, 100% { transform: rotate(45deg) translateY(0); }
    50%       { transform: rotate(45deg) translateY(4px); }
  }

  @keyframes fadeInUp {
    from { opacity: 0; transform: translateX(-50%) translateY(10px); }
    to   { opacity: 1; transform: translateX(-50%) translateY(0); }
  }

  /* ── TOOL SECTION ── */
  .tool-section {
    width: 100%;
    display: flex;
    flex-direction: column;
    align-items: center;
    padding: 60px 0 80px;
  }

  .tool-section-header {
    text-align: center;
    margin-bottom: 32px;
  }

  .tool-section-header h2 {
    font-family: 'Playfair Display', serif;
    font-size: 1.6rem;
    font-weight: 400;
    color: var(--cream);
  }

  .tool-section-header p {
    margin-top: 6px;
    font-size: 0.85rem;
    color: var(--text-muted);
  }

  /* ── MISC ── */
  .hidden { display: none !important; }

  /* ── SCROLLBAR ── */
  ::-webkit-scrollbar { width: 6px; }
  ::-webkit-scrollbar-track { background: var(--navy); }
  ::-webkit-scrollbar-thumb { background: var(--navy-light); border-radius: 3px; }
  ::-webkit-scrollbar-thumb:hover { background: var(--ink-blue); }
</style>
</head>
<body>

<!-- ═══ LANDING PAGE ═══ -->
<section class="landing" id="landingSection">
  <img src="/static/favicon.png" alt="HandyWrite" class="landing-logo" onerror="this.style.display='none'">

  <h1>Handy<em>Write</em></h1>
  <p class="landing-tagline">
    Drop in any lecture — PDF, slides, or photo.<br>
    Get back <strong>real handwritten notes</strong> in seconds.
  </p>

  <div class="features">
    <div class="feature-card">
      <span class="feature-icon">📄</span>
      <h3>Any Format</h3>
      <p>PDF, PPTX, or a photo of your slides</p>
    </div>
    <div class="feature-card">
      <span class="feature-icon">🤖</span>
      <h3>AI-Summarised</h3>
      <p>GPT distills the key ideas for you</p>
    </div>
    <div class="feature-card">
      <span class="feature-icon">✍️</span>
      <h3>Handwritten</h3>
      <p>Rendered on real lined paper, printable</p>
    </div>
  </div>

  <button class="cta-btn" id="ctaBtn">
    Try it out
    <span class="cta-arrow">→</span>
  </button>

  <div class="scroll-hint">
    <span>scroll</span>
    <div class="chevron"></div>
  </div>
</section>

<!-- ═══ TOOL SECTION ═══ -->
<section class="tool-section hidden" id="toolSection">

  <div class="tool-section-header">
    <h2>Generate Your Notes</h2>
    <p>Upload a lecture file and let HandyWrite do the rest</p>
  </div>

<!-- Main card -->
<div class="card">

  <!-- Drop zone -->
  <div class="drop-zone" id="dropZone">
    <input type="file" id="fileInput" accept=".pdf,.pptx,.png,.jpg,.jpeg">
    <div class="drop-icon-wrap">🖋️</div>
    <p><strong>Drop your lecture here</strong><br>or click to browse files</p>
    <div class="file-types">PDF · PPTX · PNG · JPG</div>
    <div class="file-name-display" id="fileNameDisplay">
      <span class="check">✓</span>
      <span id="fileName"></span>
    </div>
  </div>

  <!-- Instructions -->
  <div class="section-divider">
    <span>instructions for the AI</span>
  </div>

  <div class="instructions-wrap">
    <label for="instructions">
      Tell the AI how to take notes
      <span class="tag">optional</span>
    </label>
    <textarea id="instructions" rows="3" placeholder="e.g. focus on definitions only, skip all the math, write as if explaining to a beginner, use lots of examples..."></textarea>
  </div>

  <!-- Button -->
  <div class="btn-wrap">
    <button class="btn" id="generateBtn" disabled>
      <span class="btn-icon">✍️</span>
      Generate Handwritten Notes
    </button>
  </div>

  <!-- Loading bar -->
  <div class="progress-wrap" id="progressWrap">
    <div class="progress-bar" id="progressBar"></div>
  </div>

  <!-- Status (errors only) -->
  <div class="status" id="status"></div>

</div>

<!-- Results -->
<div class="results-header hidden" id="resultsHeader">
  <div class="results-title">
    Your Notes
    <span class="page-count" id="pageCount"></span>
  </div>
  <div style="display:flex;gap:8px;">
    <button class="dl-btn" id="dlBtn">⬇ Download PDF</button>
    <button class="dl-btn" id="printBtn">🖨 Print</button>
  </div>
</div>

<div class="pages hidden" id="pagesContainer"></div>

</section>

<script>
  // ── Landing → Tool transition ──
  const ctaBtn       = document.getElementById('ctaBtn');
  const landingSection = document.getElementById('landingSection');
  const toolSection  = document.getElementById('toolSection');

  ctaBtn.addEventListener('click', () => {
    landingSection.style.transition = 'opacity 0.4s ease, transform 0.4s ease';
    landingSection.style.opacity    = '0';
    landingSection.style.transform  = 'translateY(-20px)';
    setTimeout(() => {
      landingSection.classList.add('hidden');
      toolSection.classList.remove('hidden');
      toolSection.style.opacity   = '0';
      toolSection.style.transform = 'translateY(20px)';
      toolSection.style.transition = 'opacity 0.4s ease, transform 0.4s ease';
      requestAnimationFrame(() => requestAnimationFrame(() => {
        toolSection.style.opacity   = '1';
        toolSection.style.transform = 'translateY(0)';
        window.scrollTo({ top: 0, behavior: 'smooth' });
      }));
    }, 380);
  });
  const fileInput       = document.getElementById('fileInput');
  const dropZone        = document.getElementById('dropZone');
  const fileNameEl      = document.getElementById('fileName');
  const fileNameDisplay = document.getElementById('fileNameDisplay');
  const generateBtn     = document.getElementById('generateBtn');
  const instructions    = document.getElementById('instructions');
  const status          = document.getElementById('status');
  const progressWrap    = document.getElementById('progressWrap');
  const progressBar     = document.getElementById('progressBar');
  const pagesContainer  = document.getElementById('pagesContainer');
  const resultsHeader   = document.getElementById('resultsHeader');
  const pageCount       = document.getElementById('pageCount');
  const dlBtn           = document.getElementById('dlBtn');
  const printBtn        = document.getElementById('printBtn');

  // step → % width:  0=extract(15%), 1=AI(50%), 2=render(80%), done(100%)
  const STAGE_PROGRESS = { extracting: 15, generating: 50, rendering: 80, done: 100 };

  function setProgress(stage) {
    const pct = STAGE_PROGRESS[stage] ?? 0;
    if (pct > 0) {
      progressWrap.classList.add('visible');
      progressBar.style.width = pct + '%';
    } else {
      progressWrap.classList.remove('visible');
      progressBar.style.width = '0%';
    }
  }

  let selectedFile  = null;
  let renderedPages = [];

  // Drag & drop
  dropZone.addEventListener('dragover',  e => { e.preventDefault(); dropZone.classList.add('drag-over'); });
  dropZone.addEventListener('dragleave', e => { if (!dropZone.contains(e.relatedTarget)) dropZone.classList.remove('drag-over'); });
  dropZone.addEventListener('drop', e => {
    e.preventDefault();
    dropZone.classList.remove('drag-over');
    const f = e.dataTransfer.files[0];
    const allowedExt = ['.pdf', '.pptx', '.png', '.jpg', '.jpeg'];
    const ext = '.' + f.name.split('.').pop().toLowerCase();
    if (f && allowedExt.includes(ext)) setFile(f);
    else setStatus('Unsupported type — use PDF, PPTX, PNG, or JPG.', true);
  });

  fileInput.addEventListener('change', () => {
    if (fileInput.files[0]) setFile(fileInput.files[0]);
  });

  function setFile(f) {
    selectedFile = f;
    fileNameEl.textContent = f.name;
    fileNameDisplay.classList.add('visible');
    generateBtn.disabled = false;
    setStatus('');
    setSteps(-1);
  }

  function setStatus(msg, isError = false) {
    status.textContent = msg;
    status.className = 'status' + (isError ? ' error' : '');
  }

  function setBtnStage(stage) {
    const stages = {
      extracting: '<span class="spinner"></span> Reading lecture...',
      generating: '<span class="spinner"></span> Generating notes with AI...',
      rendering:  '<span class="spinner"></span> Rendering handwriting...',
      done:       '<span class="btn-icon">✍️</span> Generate Handwritten Notes',
    };
    generateBtn.innerHTML = stages[stage] || stages.done;
  }

  generateBtn.addEventListener('click', async () => {
    if (!selectedFile) return;
    generateBtn.disabled = true;
    setBtnStage('extracting');
    setProgress('extracting');
    pagesContainer.classList.add('hidden');
    resultsHeader.classList.add('hidden');
    renderedPages = [];

    try {
      const formData = new FormData();
      formData.append('pdf', selectedFile);
      formData.append('instructions', instructions.value.trim());

      const res  = await fetch('/generate', { method: 'POST', body: formData });
      const data = await res.json();

      if (!res.ok || data.error) throw new Error(data.error || 'Server error');

      setProgress('done');
      renderedPages = data.pages;
      showPages(data.pages);

    } catch (err) {
      setStatus('Error: ' + err.message, true);
      setProgress(null);
    } finally {
      generateBtn.disabled = false;
      setBtnStage('done');
      // fade bar out after a moment on success
      setTimeout(() => setProgress(null), 800);
    }
  });

  function showPages(pages) {
    setStatus('');

    // Results header
    pageCount.textContent = pages.length + (pages.length === 1 ? ' page' : ' pages');
    resultsHeader.classList.remove('hidden');
    dlBtn.onclick   = () => downloadAll(pages);
    printBtn.onclick = () => printPDF(pages);

    // Pages
    pagesContainer.innerHTML = '';
    pagesContainer.classList.remove('hidden');

    pages.forEach((b64, i) => {
      const wrap  = document.createElement('div');
      wrap.className = 'page-wrap';
      const img   = document.createElement('img');
      img.src     = 'data:image/png;base64,' + b64;
      img.alt     = 'Page ' + (i + 1);
      const label = document.createElement('div');
      label.className = 'page-label';
      label.textContent = 'p.' + (i + 1);
      wrap.appendChild(img);
      wrap.appendChild(label);
      pagesContainer.appendChild(wrap);
    });

    // Scroll to results smoothly
    setTimeout(() => resultsHeader.scrollIntoView({ behavior: 'smooth', block: 'start' }), 100);
  }

  async function downloadAll(pages) {
    dlBtn.textContent = '⏳ Preparing...';
    try {
      const res  = await fetch('/download', {
        method: 'POST',
        headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify({ pages })
      });
      const blob = await res.blob();
      const url  = URL.createObjectURL(blob);
      const a    = document.createElement('a');
      a.href     = url;
      a.download = 'handwritten_notes.pdf';
      a.click();
      URL.revokeObjectURL(url);
    } finally {
      dlBtn.textContent = '⬇ Download PDF';
    }
  }

  async function printPDF(pages) {
    printBtn.textContent = '⏳ Preparing...';
    try {
      const res  = await fetch('/download', {
        method: 'POST',
        headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify({ pages })
      });
      const blob = await res.blob();
      const url  = URL.createObjectURL(blob);
      // Open PDF in hidden iframe and trigger print dialog
      let iframe = document.getElementById('_printFrame');
      if (!iframe) {
        iframe = document.createElement('iframe');
        iframe.id = '_printFrame';
        iframe.style.cssText = 'position:fixed;left:-9999px;top:-9999px;width:1px;height:1px;';
        document.body.appendChild(iframe);
      }
      iframe.src = url;
      iframe.onload = () => {
        iframe.contentWindow.focus();
        iframe.contentWindow.print();
        setTimeout(() => URL.revokeObjectURL(url), 5000);
      };
    } finally {
      printBtn.textContent = '🖨 Print';
    }
  }

  // SSE progress — drives button text and loading bar through stages
  const evtSource = new EventSource('/progress');
  evtSource.onmessage = e => {
    const d = JSON.parse(e.data);
    if (d.step === 0) { setBtnStage('extracting'); setProgress('extracting'); }
    else if (d.step === 1) { setBtnStage('generating'); setProgress('generating'); }
    else if (d.step === 2) { setBtnStage('rendering');  setProgress('rendering'); }
  };
</script>
</body>
</html>"""

progress_data = {"step": -1, "msg": ""}

def set_progress(step, msg):
    progress_data["step"] = step
    progress_data["msg"]  = msg

@app.route("/")
def index():
    return render_template_string(HTML)

@app.route("/progress")
def progress():
    def stream():
        import json, time
        last = None
        for _ in range(300):
            cur = dict(progress_data)
            if cur != last and cur.get("step", -1) >= 0:
                yield f"data: {json.dumps(cur)}\n\n"
                last = dict(cur)
            time.sleep(0.3)
    from flask import Response
    return Response(stream(), mimetype="text/event-stream")

@app.route("/download", methods=["POST"])
def download():
    data      = request.get_json()
    pages_b64 = data.get("pages", [])
    images    = []
    for b64 in pages_b64:
        img_bytes = base64.b64decode(b64)
        img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        images.append(img)
    buf = io.BytesIO()
    if images:
        images[0].save(buf, format="PDF", save_all=True, append_images=images[1:])
    buf.seek(0)
    from flask import send_file
    return send_file(buf, mimetype="application/pdf",
                     as_attachment=True, download_name="handwritten_notes.pdf")

@app.route("/generate", methods=["POST"])
def generate():
    try:
        uploaded = request.files.get("pdf")
        if not uploaded:
            return jsonify({"error": "No file uploaded"}), 400

        messiness    = 0.3
        detail       = 0.5
        filename     = uploaded.filename or "upload.pdf"
        custom_instr = request.form.get("instructions", "").strip()

        set_progress(0, "Extracting content...")
        file_bytes = uploaded.read()
        raw_text   = extract_from_upload(file_bytes, filename)

        set_progress(1, "Sending to GPT...")
        notes = generate_notes(raw_text, detail=detail, custom_instructions=custom_instr)

        set_progress(2, "Rendering handwritten pages...")
        pages_b64 = render_notes_to_b64(notes, messiness=messiness)

        set_progress(3, "Done!")
        result = jsonify({"pages": pages_b64})
        set_progress(-1, "")
        return result

    except Exception as e:
        traceback.print_exc()
        set_progress(-1, "")
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    print("Starting server at http://localhost:5000")
    print(f"Font: {FONT_PATH}")
    if OPENAI_API_KEY == "YOUR_API_KEY_HERE":
        print("WARNING: Set your OpenAI API key in app.py or via OPENAI_API_KEY env var")
    port = int(os.environ.get("PORT", 5000))
    debug = os.environ.get("RAILWAY_ENVIRONMENT") is None
    app.run(host="0.0.0.0", port=port, debug=debug)